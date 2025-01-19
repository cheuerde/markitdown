# type: ignore
import base64
import binascii
import copy
import html
import json
import mimetypes
import os
import re
import shutil
import subprocess
from typing import Optional
import sys
import tempfile
import traceback
import zipfile
from xml.dom import minidom
from typing import Any, Dict, List, Optional, Union
from pathlib import Path
from urllib.parse import parse_qs, quote, unquote, urlparse, urlunparse
from warnings import warn, resetwarnings, catch_warnings

import mammoth
import markdownify
import olefile
import pandas as pd
import pdfminer
import pdfminer.high_level
import pptx

# File-format detection
import puremagic
import requests
from bs4 import BeautifulSoup
from charset_normalizer import from_path

# Optional Transcription support
IS_AUDIO_TRANSCRIPTION_CAPABLE = False
try:
    # Using warnings' catch_warnings to catch
    # pydub's warning of ffmpeg or avconv missing
    with catch_warnings(record=True) as w:
        import pydub

        if w:
            raise ModuleNotFoundError
    import speech_recognition as sr

    IS_AUDIO_TRANSCRIPTION_CAPABLE = True
except ModuleNotFoundError:
    pass
finally:
    resetwarnings()

# Optional YouTube transcription support
try:
    from youtube_transcript_api import YouTubeTranscriptApi

    IS_YOUTUBE_TRANSCRIPT_CAPABLE = True
except ModuleNotFoundError:
    pass


class _CustomMarkdownify(markdownify.MarkdownConverter):
    """
    A custom version of markdownify's MarkdownConverter. Changes include:

    - Altering the default heading style to use '#', '##', etc.
    - Removing javascript hyperlinks.
    - Truncating images with large data:uri sources.
    - Ensuring URIs are properly escaped, and do not conflict with Markdown syntax
    """

    def __init__(self, **options: Any):
        options["heading_style"] = options.get("heading_style", markdownify.ATX)
        # Explicitly cast options to the expected type if necessary
        super().__init__(**options)

    def convert_hn(self, n: int, el: Any, text: str, convert_as_inline: bool) -> str:
        """Same as usual, but be sure to start with a new line"""
        if not convert_as_inline:
            if not re.search(r"^\n", text):
                return "\n" + super().convert_hn(n, el, text, convert_as_inline)  # type: ignore

        return super().convert_hn(n, el, text, convert_as_inline)  # type: ignore

    def convert_a(self, el: Any, text: str, convert_as_inline: bool):
        """Same as usual converter, but removes Javascript links and escapes URIs."""
        prefix, suffix, text = markdownify.chomp(text)  # type: ignore
        if not text:
            return ""
        href = el.get("href")
        title = el.get("title")

        # Escape URIs and skip non-http or file schemes
        if href:
            try:
                parsed_url = urlparse(href)  # type: ignore
                if parsed_url.scheme and parsed_url.scheme.lower() not in ["http", "https", "file"]:  # type: ignore
                    return "%s%s%s" % (prefix, text, suffix)
                href = urlunparse(parsed_url._replace(path=quote(unquote(parsed_url.path))))  # type: ignore
            except ValueError:  # It's not clear if this ever gets thrown
                return "%s%s%s" % (prefix, text, suffix)

        # For the replacement see #29: text nodes underscores are escaped
        if (
            self.options["autolinks"]
            and text.replace(r"\_", "_") == href
            and not title
            and not self.options["default_title"]
        ):
            # Shortcut syntax
            return "<%s>" % href
        if self.options["default_title"] and not title:
            title = href
        title_part = ' "%s"' % title.replace('"', r"\"") if title else ""
        return (
            "%s[%s](%s%s)%s" % (prefix, text, href, title_part, suffix)
            if href
            else text
        )

    def convert_img(self, el: Any, text: str, convert_as_inline: bool) -> str:
        """Same as usual converter, but removes data URIs"""

        alt = el.attrs.get("alt", None) or ""
        src = el.attrs.get("src", None) or ""
        title = el.attrs.get("title", None) or ""
        title_part = ' "%s"' % title.replace('"', r"\"") if title else ""
        if (
            convert_as_inline
            and el.parent.name not in self.options["keep_inline_images_in"]
        ):
            return alt

        # Remove dataURIs
        if src.startswith("data:"):
            src = src.split(",")[0] + "..."

        return "![%s](%s%s)" % (alt, src, title_part)

    def convert_soup(self, soup: Any) -> str:
        return super().convert_soup(soup)  # type: ignore


class DocumentConverterResult:
    """The result of converting a document to text."""

    def __init__(self, title: Union[str, None] = None, text_content: str = ""):
        self.title: Union[str, None] = title
        self.text_content: str = text_content


class DocumentConverter:
    """Abstract superclass of all DocumentConverters."""

    def convert(
        self, local_path: str, **kwargs: Any
    ) -> Union[None, DocumentConverterResult]:
        raise NotImplementedError()


class PlainTextConverter(DocumentConverter):
    """Anything with content type text/plain"""

    def convert(
        self, local_path: str, **kwargs: Any
    ) -> Union[None, DocumentConverterResult]:
        # Guess the content type from any file extension that might be around
        content_type, _ = mimetypes.guess_type(
            "__placeholder" + kwargs.get("file_extension", "")
        )

        # Only accept text files
        if content_type is None:
            return None
        elif all(
            not content_type.lower().startswith(type_prefix)
            for type_prefix in ["text/", "application/json"]
        ):
            return None

        text_content = str(from_path(local_path).best())
        return DocumentConverterResult(
            title=None,
            text_content=text_content,
        )


class HtmlConverter(DocumentConverter):
    """Anything with content type text/html"""

    def convert(
        self, local_path: str, **kwargs: Any
    ) -> Union[None, DocumentConverterResult]:
        # Bail if not html
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".html", ".htm"]:
            return None

        result = None
        with open(local_path, "rt", encoding="utf-8") as fh:
            result = self._convert(fh.read())

        return result

    def _convert(self, html_content: str) -> Union[None, DocumentConverterResult]:
        """Helper function that converts and HTML string."""

        # Parse the string
        soup = BeautifulSoup(html_content, "html.parser")

        # Remove javascript and style blocks
        for script in soup(["script", "style"]):
            script.extract()

        # Print only the main content
        body_elm = soup.find("body")
        webpage_text = ""
        if body_elm:
            webpage_text = _CustomMarkdownify().convert_soup(body_elm)
        else:
            webpage_text = _CustomMarkdownify().convert_soup(soup)

        assert isinstance(webpage_text, str)

        return DocumentConverterResult(
            title=None if soup.title is None else soup.title.string,
            text_content=webpage_text,
        )


class RSSConverter(DocumentConverter):
    """Convert RSS / Atom type to markdown"""

    def convert(
        self, local_path: str, **kwargs
    ) -> Union[None, DocumentConverterResult]:
        # Bail if not RSS type
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".xml", ".rss", ".atom"]:
            return None
        try:
            doc = minidom.parse(local_path)
        except BaseException as _:
            return None
        result = None
        if doc.getElementsByTagName("rss"):
            # A RSS feed must have a root element of <rss>
            result = self._parse_rss_type(doc)
        elif doc.getElementsByTagName("feed"):
            root = doc.getElementsByTagName("feed")[0]
            if root.getElementsByTagName("entry"):
                # An Atom feed must have a root element of <feed> and at least one <entry>
                result = self._parse_atom_type(doc)
            else:
                return None
        else:
            # not rss or atom
            return None

        return result

    def _parse_atom_type(
        self, doc: minidom.Document
    ) -> Union[None, DocumentConverterResult]:
        """Parse the type of an Atom feed.

        Returns None if the feed type is not recognized or something goes wrong.
        """
        try:
            root = doc.getElementsByTagName("feed")[0]
            title = self._get_data_by_tag_name(root, "title")
            subtitle = self._get_data_by_tag_name(root, "subtitle")
            entries = root.getElementsByTagName("entry")
            md_text = f"# {title}\n"
            if subtitle:
                md_text += f"{subtitle}\n"
            for entry in entries:
                entry_title = self._get_data_by_tag_name(entry, "title")
                entry_summary = self._get_data_by_tag_name(entry, "summary")
                entry_updated = self._get_data_by_tag_name(entry, "updated")
                entry_content = self._get_data_by_tag_name(entry, "content")

                if entry_title:
                    md_text += f"\n## {entry_title}\n"
                if entry_updated:
                    md_text += f"Updated on: {entry_updated}\n"
                if entry_summary:
                    md_text += self._parse_content(entry_summary)
                if entry_content:
                    md_text += self._parse_content(entry_content)

            return DocumentConverterResult(
                title=title,
                text_content=md_text,
            )
        except BaseException as _:
            return None

    def _parse_rss_type(
        self, doc: minidom.Document
    ) -> Union[None, DocumentConverterResult]:
        """Parse the type of an RSS feed.

        Returns None if the feed type is not recognized or something goes wrong.
        """
        try:
            root = doc.getElementsByTagName("rss")[0]
            channel = root.getElementsByTagName("channel")
            if not channel:
                return None
            channel = channel[0]
            channel_title = self._get_data_by_tag_name(channel, "title")
            channel_description = self._get_data_by_tag_name(channel, "description")
            items = channel.getElementsByTagName("item")
            if channel_title:
                md_text = f"# {channel_title}\n"
            if channel_description:
                md_text += f"{channel_description}\n"
            if not items:
                items = []
            for item in items:
                title = self._get_data_by_tag_name(item, "title")
                description = self._get_data_by_tag_name(item, "description")
                pubDate = self._get_data_by_tag_name(item, "pubDate")
                content = self._get_data_by_tag_name(item, "content:encoded")

                if title:
                    md_text += f"\n## {title}\n"
                if pubDate:
                    md_text += f"Published on: {pubDate}\n"
                if description:
                    md_text += self._parse_content(description)
                if content:
                    md_text += self._parse_content(content)

            return DocumentConverterResult(
                title=channel_title,
                text_content=md_text,
            )
        except BaseException as _:
            print(traceback.format_exc())
            return None

    def _parse_content(self, content: str) -> str:
        """Parse the content of an RSS feed item"""
        try:
            # using bs4 because many RSS feeds have HTML-styled content
            soup = BeautifulSoup(content, "html.parser")
            return _CustomMarkdownify().convert_soup(soup)
        except BaseException as _:
            return content

    def _get_data_by_tag_name(
        self, element: minidom.Element, tag_name: str
    ) -> Union[str, None]:
        """Get data from first child element with the given tag name.
        Returns None when no such element is found.
        """
        nodes = element.getElementsByTagName(tag_name)
        if not nodes:
            return None
        fc = nodes[0].firstChild
        if fc:
            return fc.data
        return None


class WikipediaConverter(DocumentConverter):
    """Handle Wikipedia pages separately, focusing only on the main document content."""

    def convert(
        self, local_path: str, **kwargs: Any
    ) -> Union[None, DocumentConverterResult]:
        # Bail if not Wikipedia
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".html", ".htm"]:
            return None
        url = kwargs.get("url", "")
        if not re.search(r"^https?:\/\/[a-zA-Z]{2,3}\.wikipedia.org\/", url):
            return None

        # Parse the file
        soup = None
        with open(local_path, "rt", encoding="utf-8") as fh:
            soup = BeautifulSoup(fh.read(), "html.parser")

        # Remove javascript and style blocks
        for script in soup(["script", "style"]):
            script.extract()

        # Print only the main content
        body_elm = soup.find("div", {"id": "mw-content-text"})
        title_elm = soup.find("span", {"class": "mw-page-title-main"})

        webpage_text = ""
        main_title = None if soup.title is None else soup.title.string

        if body_elm:
            # What's the title
            if title_elm and len(title_elm) > 0:
                main_title = title_elm.string  # type: ignore
                assert isinstance(main_title, str)

            # Convert the page
            webpage_text = f"# {main_title}\n\n" + _CustomMarkdownify().convert_soup(
                body_elm
            )
        else:
            webpage_text = _CustomMarkdownify().convert_soup(soup)

        return DocumentConverterResult(
            title=main_title,
            text_content=webpage_text,
        )


class YouTubeConverter(DocumentConverter):
    """Handle YouTube specially, focusing on the video title, description, and transcript."""

    def convert(
        self, local_path: str, **kwargs: Any
    ) -> Union[None, DocumentConverterResult]:
        # Bail if not YouTube
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".html", ".htm"]:
            return None
        url = kwargs.get("url", "")
        if not url.startswith("https://www.youtube.com/watch?"):
            return None

        # Parse the file
        soup = None
        with open(local_path, "rt", encoding="utf-8") as fh:
            soup = BeautifulSoup(fh.read(), "html.parser")

        # Read the meta tags
        assert soup.title is not None and soup.title.string is not None
        metadata: Dict[str, str] = {"title": soup.title.string}
        for meta in soup(["meta"]):
            for a in meta.attrs:
                if a in ["itemprop", "property", "name"]:
                    metadata[meta[a]] = meta.get("content", "")
                    break

        # We can also try to read the full description. This is more prone to breaking, since it reaches into the page implementation
        try:
            for script in soup(["script"]):
                content = script.text
                if "ytInitialData" in content:
                    lines = re.split(r"\r?\n", content)
                    obj_start = lines[0].find("{")
                    obj_end = lines[0].rfind("}")
                    if obj_start >= 0 and obj_end >= 0:
                        data = json.loads(lines[0][obj_start : obj_end + 1])
                        attrdesc = self._findKey(data, "attributedDescriptionBodyText")  # type: ignore
                        if attrdesc:
                            metadata["description"] = str(attrdesc["content"])
                    break
        except Exception:
            pass

        # Start preparing the page
        webpage_text = "# YouTube\n"

        title = self._get(metadata, ["title", "og:title", "name"])  # type: ignore
        assert isinstance(title, str)

        if title:
            webpage_text += f"\n## {title}\n"

        stats = ""
        views = self._get(metadata, ["interactionCount"])  # type: ignore
        if views:
            stats += f"- **Views:** {views}\n"

        keywords = self._get(metadata, ["keywords"])  # type: ignore
        if keywords:
            stats += f"- **Keywords:** {keywords}\n"

        runtime = self._get(metadata, ["duration"])  # type: ignore
        if runtime:
            stats += f"- **Runtime:** {runtime}\n"

        if len(stats) > 0:
            webpage_text += f"\n### Video Metadata\n{stats}\n"

        description = self._get(metadata, ["description", "og:description"])  # type: ignore
        if description:
            webpage_text += f"\n### Description\n{description}\n"

        if IS_YOUTUBE_TRANSCRIPT_CAPABLE:
            transcript_text = ""
            parsed_url = urlparse(url)  # type: ignore
            params = parse_qs(parsed_url.query)  # type: ignore
            if "v" in params:
                assert isinstance(params["v"][0], str)
                video_id = str(params["v"][0])
                try:
                    youtube_transcript_languages = kwargs.get(
                        "youtube_transcript_languages", ("en",)
                    )
                    # Must be a single transcript.
                    transcript = YouTubeTranscriptApi.get_transcript(video_id, languages=youtube_transcript_languages)  # type: ignore
                    transcript_text = " ".join([part["text"] for part in transcript])  # type: ignore
                    # Alternative formatting:
                    # formatter = TextFormatter()
                    # formatter.format_transcript(transcript)
                except Exception:
                    pass
            if transcript_text:
                webpage_text += f"\n### Transcript\n{transcript_text}\n"

        title = title if title else soup.title.string
        assert isinstance(title, str)

        return DocumentConverterResult(
            title=title,
            text_content=webpage_text,
        )

    def _get(
        self,
        metadata: Dict[str, str],
        keys: List[str],
        default: Union[str, None] = None,
    ) -> Union[str, None]:
        for k in keys:
            if k in metadata:
                return metadata[k]
        return default

    def _findKey(self, json: Any, key: str) -> Union[str, None]:  # TODO: Fix json type
        if isinstance(json, list):
            for elm in json:
                ret = self._findKey(elm, key)
                if ret is not None:
                    return ret
        elif isinstance(json, dict):
            for k in json:
                if k == key:
                    return json[k]
                else:
                    ret = self._findKey(json[k], key)
                    if ret is not None:
                        return ret
        return None


class IpynbConverter(DocumentConverter):
    """Converts Jupyter Notebook (.ipynb) files to Markdown."""

    def convert(
        self, local_path: str, **kwargs: Any
    ) -> Union[None, DocumentConverterResult]:
        # Bail if not ipynb
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".ipynb":
            return None

        # Parse and convert the notebook
        result = None
        with open(local_path, "rt", encoding="utf-8") as fh:
            notebook_content = json.load(fh)
            result = self._convert(notebook_content)

        return result

    def _convert(self, notebook_content: dict) -> Union[None, DocumentConverterResult]:
        """Helper function that converts notebook JSON content to Markdown."""
        try:
            md_output = []
            title = None

            for cell in notebook_content.get("cells", []):
                cell_type = cell.get("cell_type", "")
                source_lines = cell.get("source", [])

                if cell_type == "markdown":
                    md_output.append("".join(source_lines))

                    # Extract the first # heading as title if not already found
                    if title is None:
                        for line in source_lines:
                            if line.startswith("# "):
                                title = line.lstrip("# ").strip()
                                break

                elif cell_type == "code":
                    # Code cells are wrapped in Markdown code blocks
                    md_output.append(f"```python\n{''.join(source_lines)}\n```")
                elif cell_type == "raw":
                    md_output.append(f"```\n{''.join(source_lines)}\n```")

            md_text = "\n\n".join(md_output)

            # Check for title in notebook metadata
            title = notebook_content.get("metadata", {}).get("title", title)

            return DocumentConverterResult(
                title=title,
                text_content=md_text,
            )

        except Exception as e:
            raise FileConversionException(
                f"Error converting .ipynb file: {str(e)}"
            ) from e


class BingSerpConverter(DocumentConverter):
    """
    Handle Bing results pages (only the organic search results).
    NOTE: It is better to use the Bing API
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a Bing SERP
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".html", ".htm"]:
            return None
        url = kwargs.get("url", "")
        if not re.search(r"^https://www\.bing\.com/search\?q=", url):
            return None

        # Parse the query parameters
        parsed_params = parse_qs(urlparse(url).query)
        query = parsed_params.get("q", [""])[0]

        # Parse the file
        soup = None
        with open(local_path, "rt", encoding="utf-8") as fh:
            soup = BeautifulSoup(fh.read(), "html.parser")

        # Clean up some formatting
        for tptt in soup.find_all(class_="tptt"):
            if hasattr(tptt, "string") and tptt.string:
                tptt.string += " "
        for slug in soup.find_all(class_="algoSlug_icon"):
            slug.extract()

        # Parse the algorithmic results
        _markdownify = _CustomMarkdownify()
        results = list()
        for result in soup.find_all(class_="b_algo"):
            # Rewrite redirect urls
            for a in result.find_all("a", href=True):
                parsed_href = urlparse(a["href"])
                qs = parse_qs(parsed_href.query)

                # The destination is contained in the u parameter,
                # but appears to be base64 encoded, with some prefix
                if "u" in qs:
                    u = (
                        qs["u"][0][2:].strip() + "=="
                    )  # Python 3 doesn't care about extra padding

                    try:
                        # RFC 4648 / Base64URL" variant, which uses "-" and "_"
                        a["href"] = base64.b64decode(u, altchars="-_").decode("utf-8")
                    except UnicodeDecodeError:
                        pass
                    except binascii.Error:
                        pass

            # Convert to markdown
            md_result = _markdownify.convert_soup(result).strip()
            lines = [line.strip() for line in re.split(r"\n+", md_result)]
            results.append("\n".join([line for line in lines if len(line) > 0]))

        webpage_text = (
            f"## A Bing search for '{query}' found the following results:\n\n"
            + "\n\n".join(results)
        )

        return DocumentConverterResult(
            title=None if soup.title is None else soup.title.string,
            text_content=webpage_text,
        )


class PdfConverter(DocumentConverter):
    """
    Converts PDFs to Markdown. Most style information is ignored, so the results are essentially plain-text.
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a PDF
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".pdf":
            return None

        return DocumentConverterResult(
            title=None,
            text_content=pdfminer.high_level.extract_text(local_path),
        )


class DocxConverter(HtmlConverter):
    """
    Converts DOCX files to Markdown. Style information (e.g., headings) and tables are preserved where possible.
    """

    def sanitize_filename(self, name: str) -> str:
        """Sanitizes a string to make it a valid file name across different operating systems."""
        # Normalize underscore
        name = re.sub(r"\s+", "_", name.strip())

        # Replace invalid characters with underscores
        name = re.sub(r'[\\/*?:"<>|]', "_", name)

        # Remove leading and trailing dots and spaces
        name = name.strip(" .")

        # Limit the length of the filename to a reasonable length (e.g., 251 characters)
        max_length = 251
        if len(name) > max_length:
            name = name[:max_length]

        return name

    def truncate_filename(self, name: str, max_length: int, extension: str = "") -> str:
        """Truncates the filename to ensure the final length is within the limit."""
        max_base_length = max_length - len(extension)
        if len(name) > max_base_length:
            return name[:max_base_length]
        return name

    def unique_filename(self, base_path: str, max_length: int = 251) -> str:
        """Generates a unique filename while ensuring it stays within the length limit."""
        base, ext = os.path.splitext(base_path)
        truncated_base = self.truncate_filename(base, max_length, ext)

        counter = 1
        unique_path = f"{truncated_base}{ext}"
        while os.path.exists(unique_path):
            suffix = f"_{counter}"
            # Ensure base is short enough to add the suffix
            truncated_base = self.truncate_filename(
                base, max_length - len(suffix) - len(ext)
            )
            unique_path = f"{truncated_base}{suffix}{ext}"
            counter += 1

        return unique_path

    def convert_image(self, image, output_dir: str) -> dict:
        """Handles image extraction and saving with collision avoidance and length limits."""
        os.makedirs(output_dir, exist_ok=True)

        image.alt_text = image.alt_text.replace("\n", " ")
        raw_name = image.alt_text or f"image_{hash(image)}"
        sanitized_name = self.sanitize_filename(raw_name)
        truncated_name = self.truncate_filename(sanitized_name, 251, ".png")
        image_path = os.path.join(output_dir, truncated_name + ".png")

        # Ensure unique filename
        image_path = self.unique_filename(image_path)

        try:
            with image.open() as image_bytes:
                with open(image_path, "wb") as img_file:
                    img_file.write(image_bytes.read())
            return {"src": image_path, "alt": image.alt_text}
        except Exception:
            # Return an empty src if saving fails
            return {"src": ""}

    def convert(
        self, local_path: str, **kwargs
    ) -> Union[None, DocumentConverterResult]:
        # Bail if not a DOCX
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".docx":
            return None

        try:
            with open(local_path, "rb") as docx_file:
                style_map = kwargs.get("style_map")
                image_output_dir = kwargs.get("image_output_dir", "images")

                mammoth_result = mammoth.convert_to_html(
                    docx_file,
                    style_map=style_map,
                    convert_image=mammoth.images.inline(
                        lambda img: self.convert_image(img, image_output_dir)
                    ),
                )

                html_content = mammoth_result.value
                return self._convert(html_content)
        except Exception:
            return None


class XlsxConverter(HtmlConverter):
    """
    Converts XLSX files to Markdown, with each sheet presented as a separate Markdown table.
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a XLSX
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".xlsx":
            return None

        sheets = pd.read_excel(local_path, sheet_name=None, engine="openpyxl")
        md_content = ""
        for s in sheets:
            md_content += f"## {s}\n"
            html_content = sheets[s].to_html(index=False)
            md_content += self._convert(html_content).text_content.strip() + "\n\n"

        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )


class XlsConverter(HtmlConverter):
    """
    Converts XLS files to Markdown, with each sheet presented as a separate Markdown table.
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a XLS
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".xls":
            return None

        sheets = pd.read_excel(local_path, sheet_name=None, engine="xlrd")
        md_content = ""
        for s in sheets:
            md_content += f"## {s}\n"
            html_content = sheets[s].to_html(index=False)
            md_content += self._convert(html_content).text_content.strip() + "\n\n"

        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )

class PptxConverter(HtmlConverter):
    """
    Converts PPTX files to Markdown. Supports heading, tables, images with alt text, and charts.
    Images are extracted and saved alongside the markdown file.
    """

    def sanitize_filename(self, name: str) -> str:
        """Sanitizes a string to make it a valid file name across different operating systems."""
        # Normalize underscore
        name = re.sub(r"\s+", "_", name.strip())
        # Replace invalid characters with underscores
        name = re.sub(r'[\\/*?:"<>|]', "_", name)
        # Remove leading and trailing dots and spaces
        name = name.strip(" .")
        # Limit the length of the filename
        max_length = 251
        if len(name) > max_length:
            name = name[:max_length]
        return name

    def extract_image(self, shape, output_dir: str, slide_num: int, shape_idx: int) -> Optional[dict]:
        """Extracts and saves an image from a PowerPoint shape."""
        try:
            os.makedirs(output_dir, exist_ok=True)

            # Get alt text and image data
            alt_text = ""
            try:
                alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
            except Exception:
                pass

            if not alt_text:
                alt_text = f"Slide {slide_num} - Image {shape_idx}"

            # Generate base filename
            base_name = self.sanitize_filename(f"slide_{slide_num}_image_{shape_idx}")
            
            # Check if shape actually has image data
            if not hasattr(shape, 'image') or not hasattr(shape.image, 'blob'):
                print(f"Warning: Shape in slide {slide_num} has no valid image data")
                return None

            # Determine image type and handle special formats
            try:
                image_type = shape.image.content_type.split('/')[-1].lower()
            except AttributeError:
                print(f"Warning: Cannot determine image type for shape in slide {slide_num}")
                return None

            if image_type == 'jpeg':
                image_type = 'jpg'
                final_path = os.path.join(output_dir, f"{base_name}.{image_type}")
                with open(final_path, "wb") as img_file:
                    img_file.write(shape.image.blob)
                    
            elif image_type in ['x-wmf', 'wmf']:
                # Save WMF to temporary file
                temp_wmf = os.path.join(output_dir, f"{base_name}_temp.wmf")
                final_path = os.path.join(output_dir, f"{base_name}.png")
                
                try:
                    # Write the WMF file
                    with open(temp_wmf, "wb") as img_file:
                        img_file.write(shape.image.blob)
                    
                    # Try to convert WMF to PNG using ImageMagick
                    if shutil.which('magick'):  # Check if ImageMagick is installed
                        try:
                            subprocess.run(
                                ['magick', temp_wmf, final_path], 
                                check=True,
                                capture_output=True,
                                text=True
                            )
                            os.remove(temp_wmf)  # Clean up temp file if conversion succeeded
                        except subprocess.CalledProcessError as e:
                            print(f"Warning: WMF conversion failed: {e.stderr}")
                            os.rename(temp_wmf, final_path)
                    else:
                        # If ImageMagick isn't available, just rename WMF file
                        os.rename(temp_wmf, final_path)
                        print(f"Note: ImageMagick not found. WMF file saved as {os.path.basename(final_path)}")
                except Exception as e:
                    print(f"Warning: Error handling WMF file: {str(e)}")
                    if os.path.exists(temp_wmf):
                        os.rename(temp_wmf, final_path)
                    
            else:
                # Handle all other image types normally
                final_path = os.path.join(output_dir, f"{base_name}.{image_type}")
                with open(final_path, "wb") as img_file:
                    img_file.write(shape.image.blob)

            # Verify the file was created
            if not os.path.exists(final_path):
                print(f"Warning: Failed to create image file: {os.path.basename(final_path)}")
                return None

            # Return relative path for markdown
            return {
                "src": os.path.relpath(final_path, os.path.dirname(os.path.abspath(self._current_file))),
                "alt": alt_text
            }
        except Exception as e:
            print(f"Warning: Failed to extract image: {str(e)}")
            return None


    
    def convert(self, local_path: str, **kwargs: Any) -> Union[None, DocumentConverterResult]:
        # Bail if not a PPTX
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".pptx":
            return None

        # Store the current file path for relative path calculation
        self._current_file = local_path
        
        # Create image directory next to the markdown file
        output_dir = kwargs.get("image_output_dir")
        if output_dir is None:
            output_dir = os.path.join(
                os.path.dirname(local_path),
                "images_" + os.path.basename(local_path).replace(".pptx", "")
            )

        md_content = ""
        presentation = pptx.Presentation(local_path)

        for slide_num, slide in enumerate(presentation.slides, 1):
            md_content += f"\n\n## Slide {slide_num}\n\n"

            # Handle title first if present
            if slide.shapes.title:
                md_content += f"# {slide.shapes.title.text.strip()}\n\n"

            # Process all shapes in the slide
            for shape_idx, shape in enumerate(slide.shapes, 1):
                if shape == slide.shapes.title:
                    continue  # Skip title as we've already processed it

                # Handle images
                if self._is_picture(shape):
                    image_info = self.extract_image(shape, output_dir, slide_num, shape_idx)
                    if image_info:
                        md_content += f"![{image_info['alt']}]({image_info['src']})\n\n"

                # Handle tables
                elif self._is_table(shape):
                    html_table = "<table>"
                    for row in shape.table.rows:
                        html_table += "<tr>"
                        for cell in row.cells:
                            html_table += f"<td>{html.escape(cell.text)}</td>"
                        html_table += "</tr>"
                    html_table += "</table>"
                    md_content += self._convert(html_table).text_content + "\n\n"

                # Handle charts
                elif shape.has_chart:
                    md_content += self._convert_chart_to_markdown(shape.chart) + "\n\n"

                # Handle text frames
                elif shape.has_text_frame:
                    text = shape.text.strip()
                    if text:
                        md_content += f"{text}\n\n"

            # Add slide notes if present
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    md_content += f"\n### Slide Notes\n\n{notes}\n\n"

            md_content += "---\n"  # Add separator between slides

        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip()
        )

    def _is_picture(self, shape) -> bool:
        """Check if shape is a picture."""
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER:
            if hasattr(shape, "image"):
                return True
        return False

    def _is_table(self, shape) -> bool:
        """Check if shape is a table."""
        return shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE

    def _convert_chart_to_markdown(self, chart) -> str:
        """Convert a chart to markdown table format."""
        md = "\n### Chart"
        if chart.has_title:
            md += f": {chart.chart_title.text_frame.text}"
        md += "\n\n"

        # Extract data from chart
        data = []
        try:
            # Get categories (x-axis labels)
            category_names = [c.label for c in chart.plots[0].categories]
            # Get series names (legend labels)
            series_names = [s.name for s in chart.series]
            
            # Create header row with "Category" and series names
            data.append(["Category"] + series_names)

            # Add data rows
            for idx, category in enumerate(category_names):
                row = [category]
                for series in chart.series:
                    row.append(series.values[idx])
                data.append(row)

            # Convert to markdown table
            markdown_table = []
            for row in data:
                markdown_table.append("| " + " | ".join(map(str, row)) + " |")
            
            header = markdown_table[0]
            separator = "|" + "|".join(["---"] * len(data[0])) + "|"
            
            return md + "\n".join([header, separator] + markdown_table[1:])
        except Exception as e:
            return md + f"\n\nError converting chart: {str(e)}\n"

    def _convert_html_table(self, html_table: str) -> str:
        """Convert HTML table to markdown format."""
        return self._convert(f"<html><body>{html_table}</body></html>").text_content

class MediaConverter(DocumentConverter):
    """
    Abstract class for multi-modal media (e.g., images and audio)
    """

    def _get_metadata(self, local_path):
        exiftool = shutil.which("exiftool")
        if not exiftool:
            return None
        else:
            try:
                result = subprocess.run(
                    [exiftool, "-json", local_path], capture_output=True, text=True
                ).stdout
                return json.loads(result)[0]
            except Exception:
                return None


class WavConverter(MediaConverter):
    """
    Converts WAV files to markdown via extraction of metadata (if `exiftool` is installed), and speech transcription (if `speech_recognition` is installed).
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a WAV
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".wav":
            return None

        md_content = ""

        # Add metadata
        metadata = self._get_metadata(local_path)
        if metadata:
            for f in [
                "Title",
                "Artist",
                "Author",
                "Band",
                "Album",
                "Genre",
                "Track",
                "DateTimeOriginal",
                "CreateDate",
                "Duration",
            ]:
                if f in metadata:
                    md_content += f"{f}: {metadata[f]}\n"

        # Transcribe
        if IS_AUDIO_TRANSCRIPTION_CAPABLE:
            try:
                transcript = self._transcribe_audio(local_path)
                md_content += "\n\n### Audio Transcript:\n" + (
                    "[No speech detected]" if transcript == "" else transcript
                )
            except Exception:
                md_content += (
                    "\n\n### Audio Transcript:\nError. Could not transcribe this audio."
                )

        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )

    def _transcribe_audio(self, local_path) -> str:
        recognizer = sr.Recognizer()
        with sr.AudioFile(local_path) as source:
            audio = recognizer.record(source)
            return recognizer.recognize_google(audio).strip()


class Mp3Converter(WavConverter):
    """
    Converts MP3 files to markdown via extraction of metadata (if `exiftool` is installed), and speech transcription (if `speech_recognition` AND `pydub` are installed).
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a MP3
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".mp3":
            return None

        md_content = ""

        # Add metadata
        metadata = self._get_metadata(local_path)
        if metadata:
            for f in [
                "Title",
                "Artist",
                "Author",
                "Band",
                "Album",
                "Genre",
                "Track",
                "DateTimeOriginal",
                "CreateDate",
                "Duration",
            ]:
                if f in metadata:
                    md_content += f"{f}: {metadata[f]}\n"

        # Transcribe
        if IS_AUDIO_TRANSCRIPTION_CAPABLE:
            handle, temp_path = tempfile.mkstemp(suffix=".wav")
            os.close(handle)
            try:
                sound = pydub.AudioSegment.from_mp3(local_path)
                sound.export(temp_path, format="wav")

                _args = dict()
                _args.update(kwargs)
                _args["file_extension"] = ".wav"

                try:
                    transcript = super()._transcribe_audio(temp_path).strip()
                    md_content += "\n\n### Audio Transcript:\n" + (
                        "[No speech detected]" if transcript == "" else transcript
                    )
                except Exception:
                    md_content += "\n\n### Audio Transcript:\nError. Could not transcribe this audio."

            finally:
                os.unlink(temp_path)

        # Return the result
        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )


class ImageConverter(MediaConverter):
    """
    Converts images to markdown via extraction of metadata (if `exiftool` is installed), OCR (if `easyocr` is installed), and description via a multimodal LLM (if an llm_client is configured).
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not an image
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".jpg", ".jpeg", ".png"]:
            return None

        md_content = ""

        # Add metadata
        metadata = self._get_metadata(local_path)
        if metadata:
            for f in [
                "ImageSize",
                "Title",
                "Caption",
                "Description",
                "Keywords",
                "Artist",
                "Author",
                "DateTimeOriginal",
                "CreateDate",
                "GPSPosition",
            ]:
                if f in metadata:
                    md_content += f"{f}: {metadata[f]}\n"

        # Try describing the image with GPTV
        llm_client = kwargs.get("llm_client")
        llm_model = kwargs.get("llm_model")
        if llm_client is not None and llm_model is not None:
            md_content += (
                "\n# Description:\n"
                + self._get_llm_description(
                    local_path,
                    extension,
                    llm_client,
                    llm_model,
                    prompt=kwargs.get("llm_prompt"),
                ).strip()
                + "\n"
            )

        return DocumentConverterResult(
            title=None,
            text_content=md_content,
        )

    def _get_llm_description(self, local_path, extension, client, model, prompt=None):
        if prompt is None or prompt.strip() == "":
            prompt = "Write a detailed caption for this image."

        data_uri = ""
        with open(local_path, "rb") as image_file:
            content_type, encoding = mimetypes.guess_type("_dummy" + extension)
            if content_type is None:
                content_type = "image/jpeg"
            image_base64 = base64.b64encode(image_file.read()).decode("utf-8")
            data_uri = f"data:{content_type};base64,{image_base64}"

        messages = [
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": data_uri,
                        },
                    },
                ],
            }
        ]

        response = client.chat.completions.create(model=model, messages=messages)
        return response.choices[0].message.content


class OutlookMsgConverter(DocumentConverter):
    """Converts Outlook .msg files to markdown by extracting email metadata and content.

    Uses the olefile package to parse the .msg file structure and extract:
    - Email headers (From, To, Subject)
    - Email body content
    """

    def convert(
        self, local_path: str, **kwargs: Any
    ) -> Union[None, DocumentConverterResult]:
        # Bail if not a MSG file
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".msg":
            return None

        try:
            msg = olefile.OleFileIO(local_path)
            # Extract email metadata
            md_content = "# Email Message\n\n"

            # Get headers
            headers = {
                "From": self._get_stream_data(msg, "__substg1.0_0C1F001F"),
                "To": self._get_stream_data(msg, "__substg1.0_0E04001F"),
                "Subject": self._get_stream_data(msg, "__substg1.0_0037001F"),
            }

            # Add headers to markdown
            for key, value in headers.items():
                if value:
                    md_content += f"**{key}:** {value}\n"

            md_content += "\n## Content\n\n"

            # Get email body
            body = self._get_stream_data(msg, "__substg1.0_1000001F")
            if body:
                md_content += body

            msg.close()

            return DocumentConverterResult(
                title=headers.get("Subject"), text_content=md_content.strip()
            )

        except Exception as e:
            raise FileConversionException(
                f"Could not convert MSG file '{local_path}': {str(e)}"
            )

    def _get_stream_data(
        self, msg: olefile.OleFileIO, stream_path: str
    ) -> Union[str, None]:
        """Helper to safely extract and decode stream data from the MSG file."""
        try:
            if msg.exists(stream_path):
                data = msg.openstream(stream_path).read()
                # Try UTF-16 first (common for .msg files)
                try:
                    return data.decode("utf-16-le").strip()
                except UnicodeDecodeError:
                    # Fall back to UTF-8
                    try:
                        return data.decode("utf-8").strip()
                    except UnicodeDecodeError:
                        # Last resort - ignore errors
                        return data.decode("utf-8", errors="ignore").strip()
        except Exception:
            pass
        return None


class ZipConverter(DocumentConverter):
    """Converts ZIP files to markdown by extracting and converting all contained files.

    The converter extracts the ZIP contents to a temporary directory, processes each file
    using appropriate converters based on file extensions, and then combines the results
    into a single markdown document. The temporary directory is cleaned up after processing.

    Example output format:
    ```markdown
    Content from the zip file `example.zip`:

    ## File: docs/readme.txt

    This is the content of readme.txt
    Multiple lines are preserved

    ## File: images/example.jpg

    ImageSize: 1920x1080
    DateTimeOriginal: 2024-02-15 14:30:00
    Description: A beautiful landscape photo

    ## File: data/report.xlsx

    ## Sheet1
    | Column1 | Column2 | Column3 |
    |---------|---------|---------|
    | data1   | data2   | data3   |
    | data4   | data5   | data6   |
    ```

    Key features:
    - Maintains original file structure in headings
    - Processes nested files recursively
    - Uses appropriate converters for each file type
    - Preserves formatting of converted content
    - Cleans up temporary files after processing
    """

    def convert(
        self, local_path: str, **kwargs: Any
    ) -> Union[None, DocumentConverterResult]:
        # Bail if not a ZIP
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".zip":
            return None

        # Skip Office files even though they are technically ZIP files
        office_extensions = ['.docx', '.pptx', '.xlsx']
        file_name = os.path.basename(local_path).lower()
        if any(file_name.endswith(ext) for ext in office_extensions):
            return None

        # Get parent converters list if available
        parent_converters = kwargs.get("_parent_converters", [])
        if not parent_converters:
            return DocumentConverterResult(
                title=None,
                text_content=f"[ERROR] No converters available to process zip contents from: {local_path}",
            )

        extracted_zip_folder_name = (
            f"extracted_{os.path.basename(local_path).replace('.zip', '_zip')}"
        )
        extraction_dir = os.path.normpath(
            os.path.join(os.path.dirname(local_path), extracted_zip_folder_name)
        )
        md_content = f"Content from the zip file `{os.path.basename(local_path)}`:\n\n"

        try:
            # Extract the zip file safely
            with zipfile.ZipFile(local_path, "r") as zipObj:
                # Safeguard against path traversal
                for member in zipObj.namelist():
                    member_path = os.path.normpath(os.path.join(extraction_dir, member))
                    if (
                        not os.path.commonprefix([extraction_dir, member_path])
                        == extraction_dir
                    ):
                        raise ValueError(
                            f"Path traversal detected in zip file: {member}"
                        )

                # Extract all files safely
                zipObj.extractall(path=extraction_dir)

            # Process each extracted file
            for root, dirs, files in os.walk(extraction_dir):
                for name in files:
                    file_path = os.path.join(root, name)
                    relative_path = os.path.relpath(file_path, extraction_dir)

                    # Get file extension
                    _, file_extension = os.path.splitext(name)

                    # Update kwargs for the file
                    file_kwargs = kwargs.copy()
                    file_kwargs["file_extension"] = file_extension
                    file_kwargs["_parent_converters"] = parent_converters

                    # Try converting the file using available converters
                    for converter in parent_converters:
                        # Skip the zip converter to avoid infinite recursion
                        if isinstance(converter, ZipConverter):
                            continue

                        result = converter.convert(file_path, **file_kwargs)
                        if result is not None:
                            md_content += f"\n## File: {relative_path}\n\n"
                            md_content += result.text_content + "\n\n"
                            break

            # Clean up extracted files if specified
            if kwargs.get("cleanup_extracted", True):
                shutil.rmtree(extraction_dir)

            return DocumentConverterResult(title=None, text_content=md_content.strip())

        except zipfile.BadZipFile:
            return DocumentConverterResult(
                title=None,
                text_content=f"[ERROR] Invalid or corrupted zip file: {local_path}",
            )
        except ValueError as ve:
            return DocumentConverterResult(
                title=None,
                text_content=f"[ERROR] Security error in zip file {local_path}: {str(ve)}",
            )
        except Exception as e:
            return DocumentConverterResult(
                title=None,
                text_content=f"[ERROR] Failed to process zip file {local_path}: {str(e)}",
            )


class FileConversionException(BaseException):
    pass


class UnsupportedFormatException(BaseException):
    pass


class MarkItDown:
    """(In preview) An extremely simple text-based document reader, suitable for LLM use.
    This reader will convert common file-types or webpages to Markdown."""

    def __init__(
        self,
        requests_session: Optional[requests.Session] = None,
        llm_client: Optional[Any] = None,
        llm_model: Optional[str] = None,
        style_map: Optional[str] = None,
        # Deprecated
        mlm_client: Optional[Any] = None,
        mlm_model: Optional[str] = None,
    ):
        if requests_session is None:
            self._requests_session = requests.Session()
        else:
            self._requests_session = requests_session

        # Handle deprecation notices
        #############################
        if mlm_client is not None:
            if llm_client is None:
                warn(
                    "'mlm_client' is deprecated, and was renamed 'llm_client'.",
                    DeprecationWarning,
                )
                llm_client = mlm_client
                mlm_client = None
            else:
                raise ValueError(
                    "'mlm_client' is deprecated, and was renamed 'llm_client'. Do not use both at the same time. Just use 'llm_client' instead."
                )

        if mlm_model is not None:
            if llm_model is None:
                warn(
                    "'mlm_model' is deprecated, and was renamed 'llm_model'.",
                    DeprecationWarning,
                )
                llm_model = mlm_model
                mlm_model = None
            else:
                raise ValueError(
                    "'mlm_model' is deprecated, and was renamed 'llm_model'. Do not use both at the same time. Just use 'llm_model' instead."
                )
        #############################

        self._llm_client = llm_client
        self._llm_model = llm_model
        self._style_map = style_map

        self._page_converters: List[DocumentConverter] = []

        # Register converters for successful browsing operations
        # Later registrations are tried first / take higher priority than earlier registrations
        # To this end, the most specific converters should appear below the most generic converters
        self.register_page_converter(PlainTextConverter())
        self.register_page_converter(HtmlConverter())
        self.register_page_converter(RSSConverter())
        self.register_page_converter(WikipediaConverter())
        self.register_page_converter(YouTubeConverter())
        self.register_page_converter(BingSerpConverter())
        self.register_page_converter(DocxConverter())
        self.register_page_converter(XlsxConverter())
        self.register_page_converter(XlsConverter())
        self.register_page_converter(PptxConverter())
        self.register_page_converter(WavConverter())
        self.register_page_converter(Mp3Converter())
        self.register_page_converter(ImageConverter())
        self.register_page_converter(IpynbConverter())
        self.register_page_converter(PdfConverter())
        self.register_page_converter(ZipConverter())
        self.register_page_converter(OutlookMsgConverter())

    def convert(
        self, source: Union[str, requests.Response, Path], **kwargs: Any
    ) -> DocumentConverterResult:  # TODO: deal with kwargs
        """
        Args:
            - source: can be a string representing a path either as string pathlib path object or url, or a requests.response object
            - extension: specifies the file extension to use when interpreting the file. If None, infer from source (path, uri, content-type, etc.)
        """

        # Local path or url
        if isinstance(source, str):
            if (
                source.startswith("http://")
                or source.startswith("https://")
                or source.startswith("file://")
            ):
                return self.convert_url(source, **kwargs)
            else:
                return self.convert_local(source, **kwargs)
        # Request response
        elif isinstance(source, requests.Response):
            return self.convert_response(source, **kwargs)
        elif isinstance(source, Path):
            return self.convert_local(source, **kwargs)

    def convert_local(
        self, path: Union[str, Path], **kwargs: Any
    ) -> DocumentConverterResult:  # TODO: deal with kwargs
        if isinstance(path, Path):
            path = str(path)
        # Prepare a list of extensions to try (in order of priority)
        ext = kwargs.get("file_extension")
        extensions = [ext] if ext is not None else []

        # Get extension alternatives from the path and puremagic
        base, ext = os.path.splitext(path)
        self._append_ext(extensions, ext)

        for g in self._guess_ext_magic(path):
            self._append_ext(extensions, g)

        # Convert
        return self._convert(path, extensions, **kwargs)

    # TODO what should stream's type be?
    def convert_stream(
        self, stream: Any, **kwargs: Any
    ) -> DocumentConverterResult:  # TODO: deal with kwargs
        # Prepare a list of extensions to try (in order of priority)
        ext = kwargs.get("file_extension")
        extensions = [ext] if ext is not None else []

        # Save the file locally to a temporary file. It will be deleted before this method exits
        handle, temp_path = tempfile.mkstemp()
        fh = os.fdopen(handle, "wb")
        result = None
        try:
            # Write to the temporary file
            content = stream.read()
            if isinstance(content, str):
                fh.write(content.encode("utf-8"))
            else:
                fh.write(content)
            fh.close()

            # Use puremagic to check for more extension options
            for g in self._guess_ext_magic(temp_path):
                self._append_ext(extensions, g)

            # Convert
            result = self._convert(temp_path, extensions, **kwargs)
        # Clean up
        finally:
            try:
                fh.close()
            except Exception:
                pass
            os.unlink(temp_path)

        return result

    def convert_url(
        self, url: str, **kwargs: Any
    ) -> DocumentConverterResult:  # TODO: fix kwargs type
        # Send a HTTP request to the URL
        response = self._requests_session.get(url, stream=True)
        response.raise_for_status()
        return self.convert_response(response, **kwargs)

    def convert_response(
        self, response: requests.Response, **kwargs: Any
    ) -> DocumentConverterResult:  # TODO fix kwargs type
        # Prepare a list of extensions to try (in order of priority)
        ext = kwargs.get("file_extension")
        extensions = [ext] if ext is not None else []

        # Guess from the mimetype
        content_type = response.headers.get("content-type", "").split(";")[0]
        self._append_ext(extensions, mimetypes.guess_extension(content_type))

        # Read the content disposition if there is one
        content_disposition = response.headers.get("content-disposition", "")
        m = re.search(r"filename=([^;]+)", content_disposition)
        if m:
            base, ext = os.path.splitext(m.group(1).strip("\"'"))
            self._append_ext(extensions, ext)

        # Read from the extension from the path
        base, ext = os.path.splitext(urlparse(response.url).path)
        self._append_ext(extensions, ext)

        # Save the file locally to a temporary file. It will be deleted before this method exits
        handle, temp_path = tempfile.mkstemp()
        fh = os.fdopen(handle, "wb")
        result = None
        try:
            # Download the file
            for chunk in response.iter_content(chunk_size=512):
                fh.write(chunk)
            fh.close()

            # Use puremagic to check for more extension options
            for g in self._guess_ext_magic(temp_path):
                self._append_ext(extensions, g)

            # Convert
            result = self._convert(temp_path, extensions, url=response.url, **kwargs)
        # Clean up
        finally:
            try:
                fh.close()
            except Exception:
                pass
            os.unlink(temp_path)

        return result

    def _convert(
        self, local_path: str, extensions: List[Union[str, None]], **kwargs
    ) -> DocumentConverterResult:
        error_trace = ""
        for ext in extensions + [None]:  # Try last with no extension
            for converter in self._page_converters:
                _kwargs = copy.deepcopy(kwargs)

                # Overwrite file_extension appropriately
                if ext is None:
                    if "file_extension" in _kwargs:
                        del _kwargs["file_extension"]
                else:
                    _kwargs.update({"file_extension": ext})

                # Copy any additional global options
                if "llm_client" not in _kwargs and self._llm_client is not None:
                    _kwargs["llm_client"] = self._llm_client

                if "llm_model" not in _kwargs and self._llm_model is not None:
                    _kwargs["llm_model"] = self._llm_model

                # Add the list of converters for nested processing
                _kwargs["_parent_converters"] = self._page_converters

                if "style_map" not in _kwargs and self._style_map is not None:
                    _kwargs["style_map"] = self._style_map

                # If we hit an error log it and keep trying
                try:
                    res = converter.convert(local_path, **_kwargs)
                except Exception:
                    error_trace = ("\n\n" + traceback.format_exc()).strip()

                if res is not None:
                    # Normalize the content
                    res.text_content = "\n".join(
                        [line.rstrip() for line in re.split(r"\r?\n", res.text_content)]
                    )
                    res.text_content = re.sub(r"\n{3,}", "\n\n", res.text_content)

                    # Todo
                    return res

        # If we got this far without success, report any exceptions
        if len(error_trace) > 0:
            raise FileConversionException(
                f"Could not convert '{local_path}' to Markdown. File type was recognized as {extensions}. While converting the file, the following error was encountered:\n\n{error_trace}"
            )

        # Nothing can handle it!
        raise UnsupportedFormatException(
            f"Could not convert '{local_path}' to Markdown. The formats {extensions} are not supported."
        )

    def _append_ext(self, extensions, ext):
        """Append a unique non-None, non-empty extension to a list of extensions."""
        if ext is None:
            return
        ext = ext.strip()
        if ext == "":
            return
        # if ext not in extensions:
        extensions.append(ext)

    def _guess_ext_magic(self, path):
        """Use puremagic (a Python implementation of libmagic) to guess a file's extension based on the first few bytes."""
        # Use puremagic to guess
        try:
            guesses = puremagic.magic_file(path)

            # Fix for: https://github.com/microsoft/markitdown/issues/222
            # If there are no guesses, then try again after trimming leading ASCII whitespaces.
            # ASCII whitespace characters are those byte values in the sequence b' \t\n\r\x0b\f'
            # (space, tab, newline, carriage return, vertical tab, form feed).
            if len(guesses) == 0:
                with open(path, "rb") as file:
                    while True:
                        char = file.read(1)
                        if not char:  # End of file
                            break
                        if not char.isspace():
                            file.seek(file.tell() - 1)
                            break
                    try:
                        guesses = puremagic.magic_stream(file)
                    except puremagic.main.PureError:
                        pass

            extensions = list()
            for g in guesses:
                ext = g.extension.strip()
                if len(ext) > 0:
                    if not ext.startswith("."):
                        ext = "." + ext
                    if ext not in extensions:
                        extensions.append(ext)
            return extensions
        except FileNotFoundError:
            pass
        except IsADirectoryError:
            pass
        except PermissionError:
            pass
        return []

    def register_page_converter(self, converter: DocumentConverter) -> None:
        """Register a page text converter."""
        self._page_converters.insert(0, converter)
